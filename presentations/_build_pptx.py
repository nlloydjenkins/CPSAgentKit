"""
Build copilot-studio-best-practices.pptx using the Microsoft Brand Template.
Maps each slide to the appropriate template layout and populates placeholders.
"""
import zipfile, shutil, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# ── Convert .potx → .pptx so python-pptx can open it ──────────────────────
SRC_TEMPLATE = r'c:\repos\CPSAgentKit\presentations\Microsoft_Brand_Template_May2023.potx'
WORK_TEMPLATE = r'c:\repos\CPSAgentKit\presentations\_template.pptx'
OUTPUT = r'c:\repos\CPSAgentKit\presentations\copilot-studio-best-practices-v4.pptx'

shutil.copy2(SRC_TEMPLATE, WORK_TEMPLATE)
with zipfile.ZipFile(WORK_TEMPLATE, 'r') as zin:
    ct = zin.read('[Content_Types].xml').decode('utf-8')
ct = ct.replace(
    'application/vnd.openxmlformats-officedocument.presentationml.template.main+xml',
    'application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml',
)
tmp = WORK_TEMPLATE + '.tmp'
with zipfile.ZipFile(WORK_TEMPLATE, 'r') as zin:
    with zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = ct.encode('utf-8') if item.filename == '[Content_Types].xml' else zin.read(item.filename)
            zout.writestr(item, data)
os.replace(tmp, WORK_TEMPLATE)

prs = Presentation(WORK_TEMPLATE)

# Remove all pre-existing slides from the template
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

# ── Layout index constants ─────────────────────────────────────────────────
LY_TITLE       = 0   # 1_Title_Gradient_Warm Gray  (ph0=Title, ph12=Subtitle, ph13/14=extra)
LY_SECTION     = 7   # Section_Gradient_Warm Gray   (ph0=Title, ph12=Text)
LY_AGENDA      = 11  # Agenda_1-column_Large type   (ph0=Title, ph12=Text)
LY_SUMMARY     = 15  # Summary page_Big headline    (ph0=Title, ph15=Text, ph16=Text)
LY_QUOTE       = 16  # Quote_Yellow                 (ph15=Quote, ph14=Attribution, ph13=extra)
LY_1COL        = 18  # 1-column_Text                (ph0=Title, ph15=Body)
LY_1COL_SUB    = 19  # 1-column_Text w/ subhead     (ph0=Title, ph17=Subhead, ph15=Body)
LY_2COL        = 20  # 2-column_Text                (ph0=Title, ph15=Left, ph20=Right)
LY_BLANK_HEAD  = 29  # Blank_with Head              (ph0=Title only — we add table shapes)
LY_RESOURCES   = 30  # Resources                    (ph0=Title, ph15/19/20/21=4 text areas)
LY_CLOSING     = 31  # Closing slide_with Logo      (ph0=Title)


# ── Helpers ────────────────────────────────────────────────────────────────
def _set_ph(slide, idx, text, font_size=None, bold=None, alignment=None):
    """Set placeholder text by index. Returns the text_frame."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = text
            if font_size:
                p.font.size = Pt(font_size)
            if bold is not None:
                p.font.bold = bold
            if alignment:
                p.alignment = alignment
            return tf
    return None


def _add_bullets(slide, ph_idx, items, font_size=14):
    """Add a list of bullet strings to a placeholder."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == ph_idx:
            tf = ph.text_frame
            tf.clear()
            tf.word_wrap = True
            for i, item in enumerate(items):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                # Handle bold prefix with **text**
                if '**' in item:
                    parts = item.split('**')
                    for j, part in enumerate(parts):
                        if not part:
                            continue
                        run = p.add_run()
                        run.text = part
                        run.font.size = Pt(font_size)
                        run.font.bold = (j % 2 == 1)  # odd segments are bold
                else:
                    run = p.add_run()
                    run.text = item
                    run.font.size = Pt(font_size)
                p.space_after = Pt(4)
            return tf
    return None


def _add_table(slide, rows_data, left=Inches(0.8), top=Inches(1.8),
               width=Inches(11.5), row_height=Inches(0.45), header_color=RGBColor(0x00, 0x78, 0xD4)):
    """Add a table shape to a slide. rows_data = list of lists (first row = header)."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    total_height = row_height * n_rows
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, int(total_height))
    table = table_shape.table

    # Distribute column widths evenly
    col_w = int(width / n_cols)
    for col in table.columns:
        col.width = col_w

    for r, row in enumerate(rows_data):
        for c, cell_text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = cell_text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
    return table


def _set_notes(slide, text):
    """Set speaker notes on a slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


# ── Slide 1: Title ─────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE])
_set_ph(slide, 0, 'Copilot Studio Best Practices')
_set_ph(slide, 12, 'A Practical Guide for Builders')
_set_notes(slide, 'Welcome everyone. Today we will walk through practical best practices for building agents in Microsoft Copilot Studio. This session covers platform fundamentals, agent design patterns, and the real-world gotchas you need to know before going to production.')

# ── Slide 2: Agenda ────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[LY_AGENDA])
_set_ph(slide, 0, 'Agenda')
_add_bullets(slide, 12, [
    '1.  Getting Started',
    '2.  Agent Design',
    '3.  Topics',
    '4.  Tools & Connectors',
    '5.  Multi-Agent Patterns',
    '6.  ALM & Governance',
    '7.  Gotchas & Anti-Patterns',
], font_size=18)
_set_notes(slide, 'We have seven sections to cover. We will start with platform fundamentals and key limits, move through agent design and topics, then cover tools, multi-agent patterns, ALM, and wrap up with the gotchas that catch most builders.')

# ═══════════════════════════════════════════════════════════════════════════
# PART 1 — GETTING STARTED
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[LY_SECTION])
_set_ph(slide, 0, 'Part 1: Getting Started')
_set_ph(slide, 12, 'Platform fundamentals & key limits')
_set_notes(slide, 'Let us start with the platform basics. Understanding the constraints up front saves significant rework later.')

# What is Copilot Studio?
slide = prs.slides.add_slide(prs.slide_layouts[LY_1COL])
_set_ph(slide, 0, 'What is Copilot Studio?')
_add_bullets(slide, 15, [
    'Low-code platform for building AI agents on the Power Platform',
    'Agents can answer questions, perform actions, and orchestrate workflows',
    'Publishes to Teams, web, M365 Copilot, and other channels',
    'Two orchestration modes: **Classic** (rule-based) and **Generative** (AI-planned)',
])
_set_notes(slide, 'Copilot Studio is Microsoft\'s low-code platform for building AI agents. The key distinction is between Classic and Generative orchestration. Classic gives you deterministic control with trigger phrases. Generative lets the AI plan which topics and tools to invoke based on the conversation context.')

# Key Limits — Table slide
slide = prs.slides.add_slide(prs.slide_layouts[LY_BLANK_HEAD])
_set_ph(slide, 0, 'Key Limits to Know Before You Build')
_add_table(slide, [
    ['Limit', 'Value', 'Why It Matters'],
    ['Instructions field', '8,000 characters', 'Approx. 1,500 words for your most important config'],
    ['Tools per agent (practical)', '25–30', 'Routing degrades beyond this'],
    ['Topics per agent', '1,000', '250 in Teams environments'],
    ['Connector payload', '5 MB (public) / 450 KB (GCC)', 'GCC is dramatically lower'],
    ['File knowledge sources', '500 files max', 'Does not apply to SharePoint'],
    ['SharePoint file size (no M365 Copilot)', '7 MB', 'Files over 7 MB are silently ignored'],
])
_set_notes(slide, 'These are the limits that most commonly cause issues. The 8,000 character instruction limit means you need to be concise. The 25-30 tool practical limit is where routing quality starts to degrade. And the GCC connector payload limit at 450 KB is dramatically lower than public cloud. Know these numbers before you start building.')

# Rate Limits
slide = prs.slides.add_slide(prs.slide_layouts[LY_1COL])
_set_ph(slide, 0, 'Rate Limits')
_add_bullets(slide, 15, [
    '**Trial/dev environments:** 10 requests per minute',
    '**Production (pay-as-you-go):** 100 RPM / 2,000 RPH',
    'No graceful degradation; the agent stops responding when limits are hit',
    '5-10 daily active users can trigger limits on lower tiers',
    '',
    'Check your billing tier before committing to SLAs',
])
_set_notes(slide, 'Rate limits are the number one surprise for new builders. Trial environments only allow 10 requests per minute. Even production pay-as-you-go is capped at 100 RPM. There is no graceful degradation. The agent simply stops responding. Five to ten daily active users can trigger limits on lower tiers.')

# Gen vs Classic — Table slide
slide = prs.slides.add_slide(prs.slide_layouts[LY_BLANK_HEAD])
_set_ph(slide, 0, 'Generative vs Classic Orchestration')
_add_table(slide, [
    ['', 'Classic', 'Generative'],
    ['How it routes', 'Trigger phrases', 'AI reads descriptions'],
    ['Actions', 'Called explicitly from topics', 'Invoked dynamically by the planner'],
    ['Knowledge', 'Fallback only', 'Proactively queried'],
    ['Best for', 'High-compliance, predictable flows', 'Flexible, multi-intent conversations'],
    ['Language', 'Any', 'English only (orchestration layer)'],
])
_set_notes(slide, 'This comparison helps you choose the right orchestration mode. Classic is predictable and works in any language. Generative is flexible but English-only at the orchestration layer. Most production agents use Generative for flexibility, but compliance-critical workflows should use Classic topics within a Generative agent.')

# ═══════════════════════════════════════════════════════════════════════════
# PART 2 — AGENT DESIGN
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[LY_SECTION])
_set_ph(slide, 0, 'Part 2: Agent Design')
_set_ph(slide, 12, 'Instructions, knowledge & orchestration')
_set_notes(slide, 'Now let us talk about how to design your agent effectively. The instructions field, knowledge configuration, and orchestration choices are your primary design levers.')

# Instructions
slide = prs.slides.add_slide(prs.slide_layouts[LY_2COL])
_set_ph(slide, 0, 'Writing Effective Instructions')
_add_bullets(slide, 15, [
    '**What to include:**',
    'Role and persona',
    'Response format preferences',
    'Scope boundaries',
    'Tool usage rules (use exact tool names)',
], font_size=13)
_add_bullets(slide, 20, [
    '**What to avoid:**',
    'Too vague: "Help users with their questions"',
    'Too restrictive: scripting every response',
    'Contradictory rules',
    'Referencing tools that aren\'t connected',
], font_size=13)
_set_notes(slide, 'The instructions field is your most important configuration surface. You have about 1,500 words to define your agent\'s behavior. On the left, the essential elements to include. On the right, the common mistakes. The most frequent error is referencing tools that are not actually connected to the agent.')

# T-C-R Framework
slide = prs.slides.add_slide(prs.slide_layouts[LY_1COL_SUB])
_set_ph(slide, 0, 'The T-C-R Framework for Instructions')
_set_ph(slide, 17, 'Task \u2022 Context \u2022 Requirements', font_size=16)
_add_bullets(slide, 15, [
    '**Task:** What the agent should accomplish',
    '**Context:** What information and constraints apply',
    '**Requirements:** What format, tone, and boundaries must be met',
    '',
    'Example:',
    '  "You are an IT Help Desk agent for Contoso.',
    '   You have access to the IT Knowledge Base and can create',
    '   support tickets via the CreateTicket tool.',
    '   Always verify the employee\'s department before creating tickets."',
], font_size=13)
_set_notes(slide, 'The T-C-R framework gives you a structured way to write instructions. Define the Task first, then the Context the agent operates in, then the specific Requirements for responses. The example shows a complete instruction block for an IT Help Desk agent that references a specific tool by name.')

# Three Control Layers
slide = prs.slides.add_slide(prs.slide_layouts[LY_BLANK_HEAD])
_set_ph(slide, 0, 'Three Control Layers')
_add_table(slide, [
    ['Layer', 'When to Use', 'Example'],
    ['Deterministic', 'Irreversible actions, compliance-critical', 'Payment processing, record deletion'],
    ['Hybrid (Intercept)', 'Medium-risk, needs oversight', 'Approval workflows, value-limit gates'],
    ['AI Orchestrator', 'Low-risk, flexible', 'Q&A, information lookups, multi-step research'],
], top=Inches(1.8), row_height=Inches(0.65))
_set_notes(slide, 'Not everything should be AI-orchestrated. Use deterministic control for irreversible actions like payments or record deletions. Use the hybrid intercept pattern for medium-risk actions that need approval gates. Reserve full AI orchestration for low-risk, flexible interactions like Q&A and research.')

# ═══════════════════════════════════════════════════════════════════════════
# PART 3 — TOPICS
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[LY_SECTION])
_set_ph(slide, 0, 'Part 3: Topics')
_set_ph(slide, 12, 'When and how to use them')
_set_notes(slide, 'Topics are the building blocks of agent behavior. Let us look at when and how to use them effectively in both Classic and Generative orchestration.')

# When to Use Topics
slide = prs.slides.add_slide(prs.slide_layouts[LY_2COL])
_set_ph(slide, 0, 'When to Use Topics')
_add_bullets(slide, 15, [
    '**Use topics for:**',
    'Deterministic processes with fixed steps',
    'Structured data collection (forms, intake)',
    'Compliance-critical reproducible paths',
    'Actions needing explicit user confirmation',
], font_size=13)
_add_bullets(slide, 20, [
    '**Don\'t use topics for:**',
    'Simple Q&A that knowledge can handle',
    'One-off lookups the orchestrator can route to a tool',
], font_size=13)
_set_notes(slide, 'Use topics when you need deterministic, reproducible behavior. Structured data collection, compliance-critical processes, or actions requiring explicit user confirmation. Do not create topics for simple Q&A that knowledge sources can handle or one-off lookups that the orchestrator can route directly.')

# Designing Topics for Generative Orchestration
slide = prs.slides.add_slide(prs.slide_layouts[LY_BLANK_HEAD])
_set_ph(slide, 0, 'Designing Topics for Generative Orchestration')
_add_table(slide, [
    ['Aspect', 'Classic', 'Generative'],
    ['Trigger', '5\u201310 trigger phrases', 'Clear natural language description'],
    ['Naming', 'Anything works', 'Active, descriptive: ResetPassword not Flow1'],
    ['Inputs', 'Question nodes', 'Auto-prompted from input names'],
    ['Outputs', 'Direct messages', 'Return output variables for orchestrator'],
])
_set_notes(slide, 'Topic design changes significantly between Classic and Generative orchestration. In Classic, you write 5 to 10 trigger phrases. In Generative, the AI reads the topic description to decide when to invoke it. This means descriptions become your primary routing mechanism. Name topics with active, descriptive names like ResetPassword, not Flow1.')

# Topic Design Tips
slide = prs.slides.add_slide(prs.slide_layouts[LY_1COL])
_set_ph(slide, 0, 'Topic Design Tips')
_add_bullets(slide, 15, [
    '**One topic = one intent.** Don\'t bundle "check order" and "cancel order"',
    '**Avoid overlapping descriptions.** Similar descriptions \u2192 agent invokes both',
    '**Use output variables** instead of direct messages; let the orchestrator compose the response',
    '**Define clear input parameters** with descriptions for auto-prompting',
    '',
    '**Descriptions are the primary routing mechanism.** They matter more than instructions for topic selection.',
])
_set_notes(slide, 'The most important takeaway here is that descriptions drive routing in Generative mode. One topic should map to one intent. Overlapping descriptions cause double invocation where the agent fires both topics. Use output variables instead of direct messages so the orchestrator can compose coherent responses.')

# ═══════════════════════════════════════════════════════════════════════════
# PART 4 — TOOLS & CONNECTORS
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[LY_SECTION])
_set_ph(slide, 0, 'Part 4: Tools & Connectors')
_set_ph(slide, 12, 'Connecting to external systems')
_set_notes(slide, 'Tools and connectors are how your agent interacts with external systems. Getting tool design right is critical for reliable agent behavior.')

# Tool Design Principles
slide = prs.slides.add_slide(prs.slide_layouts[LY_1COL])
_set_ph(slide, 0, 'Tool Design Principles')
_add_bullets(slide, 15, [
    '**1. Single purpose.** Each tool does one thing well',
    '**2. Deterministic.** Same inputs \u2192 same outputs',
    '**3. Names matter more than descriptions.** Use TranslateText, not Action_3',
    '**4. Curate aggressively.** Fewer high-quality tools > many overlapping tools',
    '**5. Every input needs a description.** Missing descriptions cause unnecessary user prompting',
])
_set_notes(slide, 'Five principles for tool design. Single purpose, deterministic behavior, clear naming, aggressive curation, and complete input descriptions. The naming point is critical because the model uses the tool name as a primary signal for selection. TranslateText is far better than Action_3.')

# Connector Action Inputs
slide = prs.slides.add_slide(prs.slide_layouts[LY_1COL])
_set_ph(slide, 0, 'Connector Action Input Configuration')
_add_bullets(slide, 15, [
    '**Every dynamic input needs a description.** Without one, the orchestrator prompts the user even when it has the value',
    '**State the value source:** "from the trigger context", "from step 3 output", not just the format',
    '**Lock down system fields:** Import Sequence Number, Time Zone Rule, UTC Conversion. Set to custom values or remove',
    '**Primary keys:** Set to GUID() for Dataverse "Add a new row" actions',
    '**Choice columns:** Include integer mappings in the input description, not just instructions',
])
_set_notes(slide, 'Input configuration is where most connector issues originate. Every dynamic input needs a description or the orchestrator will prompt the user unnecessarily. System fields like Import Sequence Number and Time Zone Rule should be locked down. For Dataverse actions, set primary keys to GUID values.')

# Custom Connector vs MCP Server
slide = prs.slides.add_slide(prs.slide_layouts[LY_BLANK_HEAD])
_set_ph(slide, 0, 'Custom Connector vs MCP Server')
_add_table(slide, [
    ['Factor', 'Custom Connector', 'MCP Server'],
    ['Best for', 'Power Platform-native APIs', 'External SaaS, cross-platform'],
    ['DLP', 'Full enforcement', 'Allow/deny tool controls only'],
    ['ALM', 'Included in solutions', 'Deployed separately'],
    ['Multi-platform', 'Power Platform only', 'Any MCP-aware client'],
    ['Local dev', 'Cloud only', 'Stdio transport for local testing'],
])
_set_notes(slide, 'Custom connectors are native to Power Platform with full DLP enforcement and solution-aware ALM. MCP servers are cross-platform and support local development with stdio transport. Choose custom connectors for Power Platform integration and DLP. Choose MCP for cross-platform portability or local testing.')

# Power Automate Flows
slide = prs.slides.add_slide(prs.slide_layouts[LY_1COL])
_set_ph(slide, 0, 'Power Automate Flows as Actions')
_add_bullets(slide, 15, [
    '**Flows run as the maker by default,** not the end user',
    '**100-second timeout.** If the flow takes longer, the agent receives no output',
    'Put long-running logic **after** the "Return value(s) to Copilot Studio" step',
    'Use a **dedicated service account** for production flows',
    'If the maker leaves the org, all flows using their credentials **break silently**',
])
_set_notes(slide, 'Three critical points about Power Automate flows as actions. First, they run as the maker by default, not the end user. Second, they have a hard 100-second timeout with no output if exceeded. Third, if the maker who created the connections leaves the organization, all flows break silently. Always use a dedicated service account for production flows.')

# ═══════════════════════════════════════════════════════════════════════════
# PART 5 — MULTI-AGENT PATTERNS
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[LY_SECTION])
_set_ph(slide, 0, 'Part 5: Multi-Agent Patterns')
_set_ph(slide, 12, 'Delegation, routing, and connected agents')
_set_notes(slide, 'Multi-agent architectures let you scale beyond single-agent limits, but they add complexity. Let us look at when multi-agent makes sense and the critical limitations to be aware of.')

# When to Go Multi-Agent
slide = prs.slides.add_slide(prs.slide_layouts[LY_2COL])
_set_ph(slide, 0, 'When to Go Multi-Agent')
_add_bullets(slide, 15, [
    '**Use multi-agent when:**',
    'Different domains need separate knowledge/tools',
    'Different teams maintain agents independently',
    'Single agent exceeds 25\u201330 tools or 8K chars',
], font_size=13)
_add_bullets(slide, 20, [
    '**Don\'t use multi-agent when:**',
    'A single well-designed agent covers the use case',
    'You need shared MCP tools across agents',
    'You need deterministic routing',
], font_size=13)
_set_notes(slide, 'Go multi-agent when you have distinct domains with separate knowledge bases, when different teams need to maintain agents independently, or when you exceed the 25-30 tool limit. Do not go multi-agent just because you can. A single well-designed agent is always simpler to maintain and debug.')

# MCP Limitation — Quote style
slide = prs.slides.add_slide(prs.slide_layouts[LY_QUOTE])
_set_ph(slide, 15, 'Child agents CANNOT invoke their own MCP servers.', font_size=28, bold=True)
_set_ph(slide, 14, 'All MCP calls must go through the parent agent. Configure tools on the parent; use child agents only for conversation logic and knowledge retrieval.', font_size=16)
_set_notes(slide, 'This is a critical platform limitation. When a parent agent delegates to a child agent, the child cannot invoke its own MCP servers. The MCP calls simply do not execute. All MCP tools must be configured on the parent agent directly. Use child agents only for conversation logic and knowledge retrieval.')

# Context Passing
slide = prs.slides.add_slide(prs.slide_layouts[LY_1COL])
_set_ph(slide, 0, 'Context Passing Between Agents')
_add_bullets(slide, 15, [
    'The orchestrator passes context when delegating to child agents',
    '**Connected agent responses are always summarised.** Citations and links are stripped',
    '**10-turn conversation history limit.** Store critical state in variables',
    'Test interaction flows thoroughly; receiving agents can return incomplete responses',
    '',
    'Instruct the parent to preserve child output as a labelled block rather than paraphrasing',
])
_set_notes(slide, 'Connected agent responses are always summarized by the orchestrator, which strips citations and links. The 10-turn conversation history limit means you need to store critical state in variables. Instruct the parent to preserve child output verbatim as a labelled block to prevent information loss.')

# ═══════════════════════════════════════════════════════════════════════════
# PART 6 — ALM & GOVERNANCE
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[LY_SECTION])
_set_ph(slide, 0, 'Part 6: ALM & Governance')
_set_ph(slide, 12, 'Lifecycle, security & compliance')
_set_notes(slide, 'ALM and governance are where Copilot Studio shows its biggest gaps today. Understanding these limitations helps you plan around them from the start.')

# ALM — The Honest Truth
slide = prs.slides.add_slide(prs.slide_layouts[LY_BLANK_HEAD])
_set_ph(slide, 0, 'ALM: Current State')
_add_table(slide, [
    ['Problem', 'Impact'],
    ['Managed solutions produce vague SQL errors', 'Knowledge/connection references don\'t transfer cleanly'],
    ['Deleted knowledge sources persist in the API', 'Ghost references cause import failures'],
    ['No version diffing or rollback', 'Can\'t compare versions or undo a bad publish'],
    ['Knowledge sources don\'t process on import', 'Must manually re-add in every target environment'],
], width=Inches(11.5), row_height=Inches(0.6))
_set_notes(slide, 'This table lists the current ALM pain points. Knowledge sources and connection references do not transfer cleanly between environments. There is no version diffing or rollback capability. Deleted knowledge sources persist in the API and cause import failures. These are real issues that require manual workarounds.')

# ALM — What You Should Do
slide = prs.slides.add_slide(prs.slide_layouts[LY_1COL])
_set_ph(slide, 0, 'ALM: Recommendations')
_add_bullets(slide, 15, [
    '**1. Work inside Solutions from day one**',
    '**2. Separate environments** for Dev, Test, Prod with distinct DLP policies',
    '**3. Document agent configuration manually.** Keep a changelog outside the platform',
    '**4. Test knowledge sources after every import.** Do not assume they transferred',
    '**5. Use deployment pipelines** (in-product or Azure DevOps / GitHub)',
    '',
    'Include knowledge source re-processing in your deployment runbook',
])
_set_notes(slide, 'Despite the limitations, you can manage ALM effectively. Work inside Solutions from day one. Maintain separate environments with distinct DLP policies. Document your agent configuration manually since the platform does not provide version history. Always test knowledge sources after every import.')

# Authentication
slide = prs.slides.add_slide(prs.slide_layouts[LY_BLANK_HEAD])
_set_ph(slide, 0, 'Authentication Identity Model')
_add_table(slide, [
    ['Identity', 'When to Use', 'Channel Requirement'],
    ['End user (delegated)', 'User attribution, per-user data access', 'Teams, authenticated webchat, M365 Copilot'],
    ['Service account', 'Background processing, autonomous agents', 'Any channel'],
    ['Mixed', 'Flow connections run as user or maker independently', 'Depends on each connection'],
], row_height=Inches(0.65))
_set_notes(slide, 'Choose your identity model based on data access requirements. Delegated identity gives you per-user data access but requires authenticated channels like Teams. Service accounts work for autonomous agents on any channel. Many production agents use a mixed model where flow connections run independently of the agent identity.')

# DLP
slide = prs.slides.add_slide(prs.slide_layouts[LY_1COL])
_set_ph(slide, 0, 'Data Loss Prevention (DLP)')
_add_bullets(slide, 15, [
    'DLP became **mandatory** for all tenants in early 2025',
    '',
    '**What you can control:**',
    'Block "No authentication" agents',
    'Block specific knowledge source types (SharePoint, public websites)',
    'Block specific connectors as tools',
    'Block HTTP requests',
    'Block publishing channels (Teams, Direct Line, Facebook)',
    'Control event triggers for autonomous agents',
    '',
    'If no channels are unblocked, agents cannot be published',
])
_set_notes(slide, 'DLP became mandatory for all tenants in early 2025. You can block unauthenticated agents, specific knowledge source types, connectors, HTTP requests, and publishing channels. If no channels are unblocked in your DLP policy, agents cannot be published at all. Review your DLP configuration before building.')

# ═══════════════════════════════════════════════════════════════════════════
# PART 7 — GOTCHAS & ANTI-PATTERNS
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[LY_SECTION])
_set_ph(slide, 0, 'Part 7: Gotchas & Anti-Patterns')
_set_ph(slide, 12, 'Known issues and common mistakes')
_set_notes(slide, 'Let us cover the gotchas and anti-patterns that catch most builders. These are the issues that are not obvious from the documentation but that you will encounter in real production scenarios.')

# Silent Failures
slide = prs.slides.add_slide(prs.slide_layouts[LY_BLANK_HEAD])
_set_ph(slide, 0, 'Silent Failures')
_add_table(slide, [
    ['Failure', 'What Happens'],
    ['SharePoint files > 7 MB (no M365 Copilot)', 'Silently ignored. No error, no answers'],
    ['Sensitivity-labelled documents', 'Show "Ready" but never provide responses'],
    ['Knowledge "Ready" status', 'Shows Ready, then In Progress, then Ready. First "Ready" is false'],
    ['ACS channel 28 KB limit', 'Variables silently dropped on handoff'],
    ['Deleted knowledge sources', 'Removed from UI but persist in the API'],
], row_height=Inches(0.55))
_set_notes(slide, 'Silent failures are the most frustrating aspect of the platform. SharePoint files over 7 MB are silently ignored with no error. Sensitivity-labelled documents show Ready status but never provide responses. The Knowledge Ready indicator shows Ready, then In Progress, then Ready again. The first Ready is a false positive.')

# Generative Orchestration Gotchas
slide = prs.slides.add_slide(prs.slide_layouts[LY_1COL])
_set_ph(slide, 0, 'Generative Orchestration Gotchas')
_add_bullets(slide, 15, [
    '**Instructions are guidance, not hard rules.** The agent can and will deviate. Enforce critical constraints through topic logic.',
    '**"Do not" is weaker than "always".** "Always redirect pricing to sales" beats "Do not answer pricing questions"',
    '**Long instructions dilute important rules.** Front-load your most critical constraints',
    '**Same query, different results** depending on conversation context. Test in realistic scenarios',
    '**Overlapping topic descriptions = double invocation.** Test and narrow descriptions',
])
_set_notes(slide, 'Key point: instructions are guidance, not hard rules. The agent can and will deviate. Positive framing works better than negative. Always redirect pricing to sales is more effective than Do not answer pricing questions. Front-load your most critical constraints because long instructions dilute important rules.')

# Content Filtering
slide = prs.slides.add_slide(prs.slide_layouts[LY_QUOTE])
_set_ph(slide, 15, 'Content Filtering: Zero Transparency', font_size=28, bold=True)
_set_ph(slide, 14, 'When a response is blocked: no logging, no reason code, no detail. You cannot tune or override the built-in filter. If legitimate content triggers it (medical, legal, security), your only option is a support ticket.', font_size=16)
_set_notes(slide, 'Content filtering provides zero diagnostic information when it triggers. No logging, no reason code, and no detail. You cannot tune or override the built-in filter. If legitimate content in domains like medical, legal, or security triggers it, your only recourse is a Microsoft support ticket.')

# Anti-Patterns
slide = prs.slides.add_slide(prs.slide_layouts[LY_BLANK_HEAD])
_set_ph(slide, 0, 'Anti-Patterns to Avoid')
_add_table(slide, [
    ['Anti-Pattern', 'Do This Instead'],
    ['Dumping every SharePoint site as knowledge', 'Curate. Only add sources that answer real user questions'],
    ['Using Excel/spreadsheets as knowledge', 'Use connectors for analytical data'],
    ['50+ tools on one agent', 'Split into connected agents at 25\u201330 tools'],
    ['No auth on production agents', 'Use DLP to block "No authentication"'],
    ['Testing only in the test panel', 'Test in the target channel with real user credentials'],
    ['Ignoring knowledge re-processing after import', 'Build re-processing into your deployment runbook'],
], row_height=Inches(0.5))
_set_notes(slide, 'These are the most common anti-patterns. Dumping every SharePoint site as knowledge instead of curating. Using Excel as a knowledge source when connectors are the right choice. Putting 50 or more tools on one agent instead of splitting at 25-30. And testing only in the test panel instead of the real channel.')

# Enterprise Design Rules of Thumb
slide = prs.slides.add_slide(prs.slide_layouts[LY_1COL])
_set_ph(slide, 0, 'Enterprise Design Rules of Thumb')
_add_bullets(slide, 15, [
    '**1. Agentic first,** procedural only when necessary',
    '**2. One agent = one domain** (25\u201330 tool max)',
    '**3. Topics for control,** generative for flexibility',
    '**4. Knowledge for reference,** connectors for transactions',
    '**5. Start with the narrowest scope.** Expand only when testing proves you need more',
    '**6. Name everything for the model.** If it is unclear to a human, it is unclear to the AI',
    '**7. Every write action needs confirmation**',
    '**8. DLP and ALM from day one**',
    '**9. Re-evaluate after every model upgrade**',
])
_set_notes(slide, 'Nine rules of thumb for enterprise agent design. Start agentic and go procedural only when necessary. One agent per domain with a 25-30 tool maximum. Topics for control, generative for flexibility. Name everything for the model. Every write action needs confirmation. DLP and ALM from day one. Re-evaluate after every model upgrade.')

# ═══════════════════════════════════════════════════════════════════════════
# CLOSING SLIDES
# ═══════════════════════════════════════════════════════════════════════════

# Key Takeaways — Summary layout
slide = prs.slides.add_slide(prs.slide_layouts[LY_SUMMARY])
_set_ph(slide, 0, 'Key Takeaways')
_add_bullets(slide, 15, [
    '**Instructions are your primary lever.** Get the 8,000 characters right before anything else',
    '**Know your limits.** Rate limits, tool counts, and file sizes before committing to features',
    '**Silent failures are common.** Test thoroughly and do not trust "Ready" status',
], font_size=14)
_add_bullets(slide, 16, [
    '**ALM requires discipline.** Work in solutions from day one and document everything',
    '**Start small.** Fewer tools, tighter scope, expand based on testing',
    '**Test in the real channel.** The test panel does not equal production',
], font_size=14)
_set_notes(slide, 'Six key takeaways. Instructions are your primary lever, so invest time getting the 8,000 characters right. Know your limits before committing to features. Silent failures are common, so test thoroughly. ALM requires discipline from day one. Start with a narrow scope. And always test in the real channel.')

# Resources
slide = prs.slides.add_slide(prs.slide_layouts[LY_RESOURCES])
_set_ph(slide, 0, 'Resources')
_set_ph(slide, 15, 'CPSAgentKit MCP Server\nnpx @cpsagentkit/mcp-server\nAI-assisted CPS guidance', font_size=12)
_set_ph(slide, 19, 'CPSAgentKit VS Code Extension\ngithub.com/nlloydjenkins/CPSAgentKit\nScaffolding & knowledge sync', font_size=12)
_set_ph(slide, 20, 'Microsoft Documentation\nlearn.microsoft.com/en-us/\nmicrosoft-copilot-studio/', font_size=12)
_set_ph(slide, 21, 'Power Platform Admin Center\nadmin.powerplatform.microsoft.com\nDLP, environments & monitoring', font_size=12)
_set_notes(slide, 'Four resources to continue your learning. The CPSAgentKit MCP Server provides AI-assisted Copilot Studio guidance. The VS Code extension helps with agent scaffolding and knowledge sync. Microsoft Learn has the official documentation. And the Power Platform Admin Center is where you manage DLP, environments, and monitoring.')

# Closing
slide = prs.slides.add_slide(prs.slide_layouts[LY_CLOSING])
_set_ph(slide, 0, 'Questions?')
_set_notes(slide, 'Thank you for attending. I am happy to take any questions about Copilot Studio best practices, agent design, or any of the topics we covered today.')

# ── Save ───────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f'Created {OUTPUT}')
print(f'Total slides: {len(prs.slides)}')

# Cleanup temp files
os.remove(WORK_TEMPLATE)
for f in [r'c:\repos\CPSAgentKit\presentations\_inspect_template.py']:
    if os.path.exists(f):
        os.remove(f)
